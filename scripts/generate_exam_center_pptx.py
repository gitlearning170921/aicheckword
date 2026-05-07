#!/usr/bin/env python3
"""
生成「考试训练中心」功能介绍 PPT（.pptx）。

依赖（可选）：pip install -r requirements-dev.txt
默认输出：docs/presentations/exam_center_overview.pptx
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path
from typing import Any, List, Tuple, Union


def _repo_root() -> Path:
    return Path(__file__).resolve().parents[1]


def _load_pptx():
    try:
        from pptx import Presentation  # noqa: WPS433
        from pptx.util import Pt  # noqa: WPS433

        return Presentation, Pt
    except ImportError:
        return None, None


def _set_title_font(shape: Any, size_pt: int, name: str, Pt: Any) -> None:
    if not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        p.font.name = name
        p.font.size = Pt(size_pt)


def _set_body_font(text_frame: Any, size_pt: int, name: str, Pt: Any) -> None:
    for p in text_frame.paragraphs:
        p.font.name = name
        p.font.size = Pt(size_pt)


def _add_title_slide(prs: Any, title: str, subtitle: str, font: str, Pt: Any) -> None:
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    _set_title_font(slide.shapes.title, 36, font, Pt)
    if len(slide.placeholders) > 1:
        ph = slide.placeholders[1]
        ph.text = subtitle
        _set_title_font(ph, 18, font, Pt)


def _add_bullets_slide(prs: Any, title: str, bullets: List[str], font: str, Pt: Any) -> None:
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    _set_title_font(slide.shapes.title, 28, font, Pt)
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    if not bullets:
        tf.text = "（无）"
        _set_body_font(tf, 16, font, Pt)
        return
    tf.text = bullets[0]
    for line in bullets[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0
    _set_body_font(tf, 16, font, Pt)


def build_presentation() -> Any:
    Presentation, Pt = _load_pptx()
    if Presentation is None:
        raise RuntimeError("python-pptx 未安装")

    font = "Microsoft YaHei"
    prs = Presentation()

    slides: List[Tuple[str, Union[List[str], None]]] = [
        (
            "__cover__",
            [
                "aicheckword（Quiz API）+ AI Word（考试中心页面）",
                "体考训练 · 题库扩维 · 下发与练习 · 统计闭环",
            ],
        ),
        (
            "背景与目标",
            [
                "面向医疗器械体考相关培训：支持多体考轨道（如国内/13485/MDSAP）。",
                "老师：组卷、AI 录题、复审发布、题库维护、考试下发与配置。",
                "学生：来一套练习、接收考试、作答提交、错题本与历史记录。",
                "统计：练习/考试活动汇总与筛选（aiword 本地库 + 上游快照）。",
            ],
        ),
        (
            "总体架构",
            [
                "浏览器访问 AI Word「考试训练中心」页面（Flask 模板 + exam_center.js）。",
                "页面请求统一走 /api/exam-center/*，由 aiword 代理到 aicheckword 的 /quiz/*（QUIZ_API_BASE_URL）。",
                "题库、套题、作答、阅卷等核心状态在 aicheckword（FastAPI + MySQL + Chroma 等）。",
                "aiword 侧 SQLite：活动记录、ingest/review 任务与上游 job 快照等，便于列表与轮询展示。",
            ],
        ),
        (
            "老师端能力（概要）",
            [
                "套题：生成、列表、详情、删除；AI 复审与发布。",
                "题库：AI 批量录题（异步 job + 轮询）、单题编辑/删除、检索与筛选。",
                "考试：组卷下发、本地任务列表、与上游 assignments 协同（视部署配置）。",
                "配置：考试与录题相关系统设置；法规/标准更新提示（长耗时接口）。",
                "项目案例：已训练案例列表下拉（与上游 project-cases 一致）。",
            ],
        ),
        (
            "学生端能力（概要）",
            [
                "来一套（练习）：按体考类型、考试类型、难度、题量生成练习套题并作答。",
                "考试任务：老师下发后的列表、开考、与练习类似的作答与提交流程。",
                "错题本、未练题量等辅助入口（调用上游 quiz 能力）。",
                "前端：web/static/js/exam_center.js 渲染题干与选项/主观题输入。",
            ],
        ),
        (
            "统计端能力（概要）",
            [
                "看板与记录列表：练习/考试次数、通过情况、筛选与刷新。",
                "数据来自 aiword 聚合接口 + 上游返回字段（以实际部署为准）。",
            ],
        ),
        (
            "上游 Quiz API（aicheckword 摘要）",
            [
                "组卷与练习：/quiz/sets/generate、/quiz/practice/generate-set、/quiz/practice/submit。",
                "录题：/quiz/bank/ingest-by-ai 与 ingest job 查询。",
                "题库与套题：bank/questions、sets CRUD、publish、review-by-ai。",
                "作答与阅卷：attempts、grading-status、auto-grade / cache 等（按题型与规则）。",
                "工具：/quiz/tools/project-cases 等项目案例列表。",
            ],
        ),
        (
            "考试类型与项目案例",
            [
                "考试类型（与体考轨道正交）：daily（日常）、new_standard（新标发布）、project_case（项目案例）。",
                "project_case 须选择已训练入库的 project_cases.id，并与向量检索 case_id 对齐。",
                "命题素材：项目案例类从主库 category=project_case + case_id 取证（见 quiz/service 与 knowledge_base）。",
            ],
        ),
        (
            "踩坑：接口与参数",
            [
                "直连 aicheckword 时：若请求体仅含 examCategory（驼峰），旧版 Pydantic 可能未映射到 exam_category，",
                "  且 QuizIngestByAIRequest 使用 extra=ignore，会导致考试类型回退为 daily，录题走混合知识源。",
                "修复：Quiz 相关请求模型为 exam_category 增加 validation_alias（examCategory / exam_category）。",
                "aiword 代理层 _expand_quiz_request_body 会同步写入多种键名；前端亦建议蛇形/驼峰双写关键字段。",
            ],
        ),
        (
            "踩坑：命题与展示",
            [
                "项目案例：练习套题常混用题库缓存 + 补题生成，AI 录题则几乎全量现生成，体感不一致时需核对取证与 top_k。",
                "题干 stem：模型易把 evidence 大段原文粘进 stem（各题型），与选项/答案混淆。",
                "处理：prompt 中 5a 条明确 stem 与 evidence 分离；落库与读库时 _trim_embedded_evidence_from_stem 剥离重复段。",
                "前端：长题干需 pre-wrap、换行与安全转义；独立「命题材料」区块曾加后按产品要求已移除。",
            ],
        ),
        (
            "运维与配置（提要）",
            [
                "aiword：QUIZ_API_BASE_URL、超时、集成密钥（请求头）等需在环境/系统配置中正确设置。",
                "aicheckword：向量库（Chroma）与 MySQL 题库表一致；录题为后台线程，注意上游可用性与超时。",
            ],
        ),
        (
            "扩展方向与维护",
            [
                "阅卷规则与主观题策略细化；更多题型与知识点权重可视化。",
                "权限与审计：老师/学生/统计端 gate 与操作留痕。",
                "更新本 PPT：编辑本脚本中 slides 列表后重新运行生成命令。",
            ],
        ),
    ]

    for spec in slides:
        t, body = spec
        if t == "__cover__":
            assert body is not None
            sub = body[0] + "\n" + body[1]
            _add_title_slide(prs, "考试训练中心 — 功能与实现提要", sub, font, Pt)
        else:
            _add_bullets_slide(prs, t, body or [], font, Pt)

    return prs


def main() -> int:
    parser = argparse.ArgumentParser(description="生成考试中心介绍 PPT")
    parser.add_argument(
        "--out",
        type=Path,
        default=None,
        help="输出 .pptx 路径（默认：仓库 docs/presentations/exam_center_overview.pptx）",
    )
    args = parser.parse_args()
    root = _repo_root()
    out = args.out if args.out is not None else root / "docs" / "presentations" / "exam_center_overview.pptx"

    if _load_pptx()[0] is None:
        print("缺少依赖 python-pptx。请执行：pip install -r requirements-dev.txt", file=sys.stderr)
        return 1

    out.parent.mkdir(parents=True, exist_ok=True)
    prs = build_presentation()
    prs.save(str(out))
    print(f"已生成: {out}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
